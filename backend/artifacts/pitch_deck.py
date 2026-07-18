"""
artifacts/pitch_deck.py
Generates a premium, dark technical-themed .pptx pitch deck from the pipeline result JSON.

Design principles:
  • High-quality dark cybertech background images with grids and circuits
  • Dynamic brand colours extracted from the UI designer's design_system
  • Cards styled in deep dark panels (Slate-900) with light slate readable text
  • High-contrast glowing neon typography for Slide 11 (Score) and Slide 13 (CTA)
  • S-curve workflow flowchart arrows on User Journey (Slide 5)
  • Database schema connection vector lines on Tech Architecture (Slide 7)
  • Timeline arrow vector on Roadmap (Slide 12)
  • UI mockup illustration on the Design slide
  • Native Column Chart on the Score slide
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ── Paths to slide background resources ───────────────────────────────────────
RESOURCES_DIR = Path(__file__).parent / "resources"
COVER_BG_PATH = str(RESOURCES_DIR / "cover_bg.jpg")
SLIDE_BG_PATH = str(RESOURCES_DIR / "slide_bg.jpg")

# ── Dark Tech Theme Brand Palette ─────────────────────────────────────────────
PRIMARY     = RGBColor(0x63, 0x66, 0xF1)   # Indigo-500
ACCENT      = RGBColor(0x10, 0xB9, 0x81)   # Emerald-500
ACCENT2     = RGBColor(0xF5, 0x9E, 0x0B)   # Amber-500
DANGER      = RGBColor(0xEF, 0x44, 0x44)   # Red-500
DARK_BG     = RGBColor(0x0A, 0x0F, 0x1E)   # Deep Navy (Fallback)
CARD_BG     = RGBColor(0x11, 0x18, 0x27)   # Slate-900 (High Contrast Dark Slate Card)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # Pure White
TEXT_TITLE  = RGBColor(0xFF, 0xFF, 0xFF)   # White Titles
TEXT_BODY   = RGBColor(0xCB, 0xD5, 0xE1)   # Slate-300 (Readable Light Slate Body)
TEXT_MUTED  = RGBColor(0x94, 0xA3, 0xB8)   # Slate-400 (Muted Slate Grey)

# ── High Contrast Neon Colors for Hud Indicators ──────────────────────────────
NEON_CYAN   = RGBColor(0x00, 0xF5, 0xD4)   # Glowing Cyan
NEON_GREEN  = RGBColor(0x39, 0xFF, 0x14)   # Glowing Tech Neon Green

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Colour helpers
# ─────────────────────────────────────────────────────────────────────────────

def _parse_hex(hex_str: str, default: RGBColor) -> RGBColor:
    """Parse '#RRGGBB' → RGBColor. Returns default on any error."""
    try:
        h = hex_str.lstrip("#")
        if len(h) != 6:
            return default
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, AttributeError, TypeError):
        return default


def _mix(c1: RGBColor, c2: RGBColor, t: float = 0.5) -> RGBColor:
    """Linear interpolation between two RGBColors (t=0→c1, t=1→c2)."""
    return RGBColor(
        round(c1[0] + (c2[0] - c1[0]) * t),
        round(c1[1] + (c2[1] - c1[1]) * t),
        round(c1[2] + (c2[2] - c1[2]) * t),
    )


def _score_color(val: float, primary: RGBColor, accent: RGBColor) -> RGBColor:
    """Returns dynamic neon green/cyan/red rating colors to prevent dark values."""
    if val >= 7.5:
        return NEON_GREEN
    if val >= 5.0:
        return NEON_CYAN
    return RGBColor(0xFF, 0x55, 0x55)  # Glowing Tech Red


# ─────────────────────────────────────────────────────────────────────────────
# Vector Graphic Primitives (Arrows & Lines)
# ─────────────────────────────────────────────────────────────────────────────

def _draw_right_arrow(slide: Any, left: float, top: float, length: float, color: RGBColor) -> None:
    """Draw a horizontal arrow pointing to the right."""
    # Shaft
    _rect(slide, left, top + 0.08, length - 0.15, 0.04, color)
    # Head
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(left + length - 0.15), Inches(top), Inches(0.18), Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 90


def _draw_left_arrow(slide: Any, left: float, top: float, length: float, color: RGBColor) -> None:
    """Draw a horizontal arrow pointing to the left."""
    # Shaft
    _rect(slide, left + 0.15, top + 0.08, length - 0.15, 0.04, color)
    # Head
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(left), Inches(top), Inches(0.18), Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 270


def _draw_down_arrow(slide: Any, left: float, top: float, length: float, color: RGBColor) -> None:
    """Draw a vertical arrow pointing down."""
    # Shaft
    _rect(slide, left + 0.08, top, 0.04, length - 0.15, color)
    # Head
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(left), Inches(top + length - 0.15), Inches(0.20), Inches(0.18)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 180


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def generate_pitch_deck(pipeline_result: dict[str, Any]) -> bytes:
    """Build a polished dark technical-themed .pptx from the pipeline JSON and return raw bytes."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    idea    = pipeline_result.get("idea", "Your Startup")
    domain  = pipeline_result.get("domain", "SaaS")
    r1      = pipeline_result.get("round1", {})
    pm      = r1.get("pm", {})
    ui      = r1.get("ui", {})
    backend = r1.get("backend", {})
    mktg    = r1.get("marketing", {})
    r2      = pipeline_result.get("round2_critique", {})
    score   = pipeline_result.get("startup_score", {})

    # Extract brand colours and font family
    ds      = ui.get("design_system") or ui.get("visual_style", {})
    primary = _parse_hex(ds.get("primary_color", ""), PRIMARY)
    accent  = _parse_hex(ds.get("accent_color",  ""), ACCENT)
    font    = ds.get("font", "Montserrat")
    card_bg = _mix(CARD_BG, primary, 0.12)

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    _slide_cover(prs, idea, pm.get("tagline", ""), domain, primary, accent, font)

    # ── SLIDE 2: Problem ──────────────────────────────────────────────────────
    _slide_text_focus(
        prs,
        title="The Problem",
        icon="🚨",
        body=pm.get("problem", ""),
        sub=f'Target Audience: {", ".join(pm.get("target_users", []) or [pm.get("target_user", "")])}',
        primary=primary, accent=primary, font_name=font,
    )

    # ── SLIDE 3: Solution ─────────────────────────────────────────────────────
    _slide_text_focus(
        prs,
        title="Our Solution",
        icon="💡",
        body=pm.get("solution", ""),
        sub=pm.get("mvp_scope", ""),
        primary=primary, accent=accent, font_name=font,
    )

    # ── SLIDE 4: MVP Features ─────────────────────────────────────────────────
    _slide_bullets(
        prs,
        title="MVP Features",
        icon="🛠",
        bullets=pm.get("mvp_features", []),
        primary=primary, accent=primary,
        subtitle=f'Timeline: {pm.get("timeline_months", 4)} months to launch',
        font_name=font,
    )

    # ── SLIDE 5: User Journey (Flowchart Timeline!) ───────────────────────────
    _slide_numbered_steps(
        prs,
        title="User Journey Flow",
        steps=pm.get("user_journey", []),
        primary=primary, accent=accent, font_name=font,
    )

    # ── SLIDE 6: UI & Design Vision (Browser Mockup Graphic!) ─────────────────
    screens      = [s.get("name", "") for s in ui.get("screens", []) if isinstance(s, dict)]
    design_style = ds.get("style", "")
    _slide_ui_mockup(
        prs,
        title="Product Design",
        screens=screens[:4],
        style_desc=design_style,
        design_goal=ui.get('design_goal', 'Clean dashboard interface'),
        primary=primary, accent=accent, card_bg=card_bg, font_name=font,
    )

    # ── SLIDE 7: Technical Architecture (Schema Diagram!) ─────────────────────
    tech        = backend.get("tech_choices", {})
    tech_items  = [f"{k.capitalize()}: {v}" for k, v in tech.items()]
    stack       = backend.get("tech_stack", [])
    entities    = backend.get("entities", [])
    _slide_tech_schema(
        prs,
        title="Technical Architecture",
        stack_items=[f"▸  {t}" for t in (stack or tech_items)[:6]],
        entities=entities[:2],
        primary=primary, accent=accent, card_bg=card_bg, font_name=font,
    )

    # ── SLIDE 8: Market Opportunity ───────────────────────────────────────────
    tam_sam_som = mktg.get("tam_sam_som", {})
    market_items= [f"{k}: {v}" for k, v in tam_sam_som.items()]
    _slide_two_col(
        prs,
        title="Market Opportunity",
        left_header="Market Size",
        left_items=[f"📊  {m}" for m in market_items] or ["📊  Large and growing market"],
        right_header="Competitive Edge",
        right_items=[
            f"🎯  {mktg.get('positioning', '')}",
        ] + [f"◆  {c}" for c in mktg.get("competitors", [])[:4]],
        primary=primary, accent=accent, card_bg=card_bg, font_name=font,
    )

    # ── SLIDE 9: Go-To-Market ─────────────────────────────────────────────────
    channels = mktg.get("channels", [])
    pricing  = mktg.get("pricing", {})
    price_lines = [f"{tier.capitalize()}: {desc}" for tier, desc in pricing.items()] if isinstance(pricing, dict) else []
    _slide_two_col(
        prs,
        title="Go-To-Market Strategy",
        left_header="Acquisition Channels",
        left_items=[f"→  {c}" for c in channels[:6]],
        right_header="Pricing Tiers",
        right_items=[f"▸  {p}" for p in (price_lines or [
            f"CAC: {mktg.get('target_cac', 'TBD')}",
            f"LTV: {mktg.get('target_ltv', 'TBD')}",
        ])[:5]],
        primary=primary, accent=accent, card_bg=card_bg, font_name=font,
    )

    # ── SLIDE 10: Risks & Mitigations ─────────────────────────────────────────
    concerns   = (r2.get("investor_concerns", []) + r2.get("skeptic_flags", []))[:4]
    revisions  = r2.get("revisions", {}).get("pm", {})
    diff_text  = revisions.get("differentiation", "")
    rev_scope  = revisions.get("revised_mvp_scope", "")
    _slide_two_col(
        prs,
        title="Validation — Risks & Mitigations",
        left_header="Key Risks Identified",
        left_items=[f"⚠  {c}" for c in concerns],
        right_header="Our Response",
        right_items=(
            [f"✓  {diff_text}"] if diff_text else []
        ) + (
            [f"✓  {rev_scope}"] if rev_scope else []
        ) + [f"✓  {r}" for r in revisions.get("updated_roadmap", [])[:3]],
        primary=primary, accent=accent, card_bg=card_bg, font_name=font,
    )

    # ── SLIDE 11: Startup Score (Native Column Chart!) ─────────────────────────
    _slide_score(prs, score, primary=primary, accent=accent, card_bg=card_bg, font_name=font)

    # ── SLIDE 12: Roadmap (Timeline Diagram Infographic!) ──────────────────────
    roadmap = pm.get("roadmap", []) or revisions.get("updated_roadmap", [])
    _slide_roadmap_timeline(
        prs,
        title="Roadmap Timeline",
        roadmap=roadmap[:3],
        primary=primary, accent=accent, card_bg=card_bg, font_name=font,
    )

    # ── SLIDE 13: CTA / End card (With Neon Cyan / Green glow) ─────────────────
    _slide_cover(
        prs,
        idea="Let's build the right thing.",
        tagline=f'"{pm.get("tagline", idea)}"',
        domain="Validated by 6 AI Co-Founders",
        primary=primary, accent=accent,
        is_end_card=True, font_name=font
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# Core slide primitives
# ─────────────────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _fill_bg(slide: Any, color: RGBColor = DARK_BG, bg_path: str | None = None) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    if bg_path and Path(bg_path).exists():
        try:
            slide.shapes.add_picture(bg_path, 0, 0, SLIDE_W, SLIDE_H)
        except Exception as e:
            print(f"[pitch_deck] Failed to add background image {bg_path}: {e}")


def _rect(
    slide: Any,
    left: float, top: float, width: float, height: float,
    color: RGBColor,
    border: bool = False,
    border_color: RGBColor | None = None,
) -> Any:
    """Add a solid filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border_color or color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _txt(
    slide: Any,
    text: str,
    left: float, top: float, width: float, height: float,
    size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = TEXT_BODY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
    font_name: str = "Segoe UI",
) -> None:
    """Add a textbox with a single formatted run."""
    if not text:
        return
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _multiline_txt(
    slide: Any,
    lines: list[str],
    left: float, top: float, width: float, height: float,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_BODY,
    line_spacing_pt: float = 4.0,
    font_name: str = "Segoe UI",
) -> None:
    """Add a textbox with multiple paragraphs, one per line."""
    if not lines:
        return
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = str(line)
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Slide templates (Dark Tech Theme)
# ─────────────────────────────────────────────────────────────────────────────

def _slide_cover(
    prs: Presentation,
    idea: str,
    tagline: str,
    domain: str,
    primary: RGBColor = PRIMARY,
    accent: RGBColor = ACCENT,
    font_name: str = "Segoe UI",
    is_end_card: bool = False,
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, COVER_BG_PATH)

    # Left accent bar (full height)
    _rect(slide, 0, 0, 0.14, 7.5, primary)

    # Glass container panel overlay to mask circuit lines and protect text visibility
    _rect(slide, 0.45, 1.4, 10.0, 4.4, CARD_BG)
    _rect(slide, 0.45, 1.4, 0.1, 4.4, NEON_CYAN)

    # Domain / label badge
    if domain:
        _rect(slide, 0.8, 0.5, 3.0, 0.45, _mix(DARK_BG, accent, 0.2))
        _txt(slide, f"  {domain.upper()}  ", 0.8, 0.52, 3.0, 0.4,
             size=11, bold=True, color=NEON_GREEN, font_name=font_name)

    # Main title - Glowing Neon Cyan title inside panel
    _txt(slide, idea, 0.8, 1.7, 9.2, 2.0,
         size=38 if is_end_card else 42,
         bold=True, color=NEON_CYAN, font_name=font_name)

    # Divider inside panel
    _rect(slide, 0.8, 3.9, 8.0, 0.04, _mix(DARK_BG, primary, 0.4))

    # Tagline - Pure White text on dark card
    if tagline:
        _txt(slide, tagline, 0.8, 4.1, 9.2, 1.4,
             size=18, italic=is_end_card, color=WHITE, font_name=font_name)

    # Footer label - Neon Green footer
    footer_color = NEON_GREEN if is_end_card else TEXT_MUTED
    footer = "Daedalus-AI" if is_end_card else "Daedalus--AI"
    _txt(slide, footer, 0.8, 6.9, 9.0, 0.5,
         size=11, color=footer_color, font_name=font_name)


def _slide_text_focus(
    prs: Presentation,
    title: str,
    icon: str,
    body: str,
    sub: str,
    primary: RGBColor,
    accent: RGBColor,
    font_name: str = "Segoe UI",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    # Left accent strip
    _rect(slide, 0, 0, 0.1, 7.5, primary)

    # Title
    _txt(slide, f"{icon}  {title}", 0.35, 0.3, 12.5, 0.9,
         size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    # Large body text in a card
    _rect(slide, 0.35, 1.55, 12.55, 4.2, CARD_BG)
    _txt(slide, body, 0.65, 1.75, 12.0, 3.8,
         size=18, color=TEXT_BODY, wrap=True, font_name=font_name)

    # Sub-text card
    if sub:
        _rect(slide, 0.35, 5.9, 12.55, 1.35, _mix(DARK_BG, accent, 0.12))
        _txt(slide, sub, 0.65, 6.0, 12.0, 1.1,
             size=14, italic=True, color=_mix(TEXT_BODY, accent, 0.4), font_name=font_name)


def _slide_bullets(
    prs: Presentation,
    title: str,
    icon: str,
    bullets: list[str],
    primary: RGBColor,
    accent: RGBColor,
    subtitle: str = "",
    font_name: str = "Segoe UI",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    _rect(slide, 0, 0, 0.1, 7.5, primary)
    _txt(slide, f"{icon}  {title}", 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    if subtitle:
        _txt(slide, subtitle, 0.35, 1.45, 12.5, 0.55, size=14, italic=True, color=TEXT_MUTED, font_name=font_name)

    y0 = 2.1 if subtitle else 1.65
    step = 0.68
    for i, bullet in enumerate(bullets[:8]):
        y = y0 + i * step
        # Accent marker dot
        _rect(slide, 0.35, y + 0.2, 0.12, 0.28, accent)
        # Bullet text
        _txt(slide, str(bullet), 0.6, y, 12.3, 0.62, size=16, color=TEXT_BODY, font_name=font_name)


def _slide_numbered_steps(
    prs: Presentation,
    title: str,
    steps: list[str],
    primary: RGBColor,
    accent: RGBColor,
    font_name: str = "Segoe UI",
) -> None:
    """User Journey Slide — styled as an active S-curve workflow layout with connector arrows."""
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    _rect(slide, 0, 0, 0.1, 7.5, primary)
    _txt(slide, f"🔢  {title}", 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    display_steps = steps[:6]
    # S-curve flowchart coordinate mapping (3 columns, 2 rows)
    card_w = 4.1
    card_h = 2.0
    coords = [
        (0.35, 1.6),   # Step 1
        (4.6, 1.6),    # Step 2
        (8.85, 1.6),   # Step 3
        (8.85, 4.4),   # Step 4
        (4.6, 4.4),    # Step 5
        (0.35, 4.4)    # Step 6
    ]

    for i, step in enumerate(display_steps):
        if i >= len(coords):
            break
        x, y = coords[i]

        # Draw step card
        _rect(slide, x, y, card_w, card_h, CARD_BG)
        # Left accent number block
        _rect(slide, x, y, 0.45, card_h, _mix(DARK_BG, primary, 0.4))
        _txt(slide, str(i + 1), x + 0.05, y + 0.65, 0.35, 0.55,
             size=18, bold=True, color=NEON_CYAN, align=PP_ALIGN.CENTER, font_name=font_name)
        # Step text
        _txt(slide, str(step), x + 0.55, y + 0.15, card_w - 0.7, card_h - 0.3,
             size=13, color=TEXT_BODY, font_name=font_name)

    # ── Draw Workflow Connection Arrows ───────────────────────────────────────
    if len(display_steps) > 1:
        # Step 1 → Step 2 (Right)
        _draw_right_arrow(slide, 4.45, 2.45, 0.15, accent)
    if len(display_steps) > 2:
        # Step 2 → Step 3 (Right)
        _draw_right_arrow(slide, 8.7, 2.45, 0.15, accent)
    if len(display_steps) > 3:
        # Step 3 → Step 4 (Down)
        _draw_down_arrow(slide, 10.9, 3.65, 0.75, accent)
    if len(display_steps) > 4:
        # Step 4 → Step 5 (Left)
        _draw_left_arrow(slide, 8.7, 5.25, 0.15, accent)
    if len(display_steps) > 5:
        # Step 5 → Step 6 (Left)
        _draw_left_arrow(slide, 4.45, 5.25, 0.15, accent)


def _slide_two_col(
    prs: Presentation,
    title: str,
    left_header: str,
    left_items: list[str],
    right_header: str,
    right_items: list[str],
    primary: RGBColor,
    accent: RGBColor,
    card_bg: RGBColor,
    font_name: str = "Segoe UI",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    _rect(slide, 0, 0, 0.1, 7.5, primary)
    _txt(slide, title, 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    col_w = 6.0
    col_h = 5.5
    x_l = 0.35
    x_r = 6.95
    y0  = 1.5

    # Left card
    _rect(slide, x_l, y0, col_w, col_h, card_bg)
    _rect(slide, x_l, y0, col_w, 0.55, _mix(DARK_BG, accent, 0.3))
    _txt(slide, left_header, x_l + 0.2, y0 + 0.1, col_w - 0.3, 0.45,
         size=15, bold=True, color=accent, font_name=font_name)
    _multiline_txt(slide, left_items[:7],
                   x_l + 0.2, y0 + 0.7, col_w - 0.35, col_h - 0.9,
                   size=14, color=TEXT_BODY, font_name=font_name)

    # Right card
    _rect(slide, x_r, y0, col_w, col_h, card_bg)
    _rect(slide, x_r, y0, col_w, 0.55, _mix(DARK_BG, primary, 0.3))
    _txt(slide, right_header, x_r + 0.2, y0 + 0.1, col_w - 0.3, 0.45,
         size=15, bold=True, color=primary, font_name=font_name)
    _multiline_txt(slide, right_items[:7],
                   x_r + 0.2, y0 + 0.7, col_w - 0.35, col_h - 0.9,
                   size=14, color=TEXT_BODY, font_name=font_name)


# ─────────────────────────────────────────────────────────────────────────────
# Visual Graphic Slides (UI Mockup, Schema, Charts, Roadmap Timelines)
# ─────────────────────────────────────────────────────────────────────────────

def _slide_ui_mockup(
    prs: Presentation,
    title: str,
    screens: list[str],
    style_desc: str,
    design_goal: str,
    primary: RGBColor,
    accent: RGBColor,
    card_bg: RGBColor,
    font_name: str = "Outfit",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    _rect(slide, 0, 0, 0.1, 7.5, primary)
    _txt(slide, f"🎨  {title}", 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    # Left: Text specs
    _rect(slide, 0.35, 1.6, 5.8, 5.4, card_bg)
    _txt(slide, "Design Language System", 0.55, 1.8, 5.4, 0.4, size=16, bold=True, color=accent, font_name=font_name)
    _multiline_txt(
        slide,
        [
            f"Style Guide: {style_desc or 'Modern dark glassmorphism'}",
            f"Fonts: {font_name}, Sans-Serif",
            f"Brand Colors: {accent[0]}R {accent[1]}G {accent[2]}B (Primary)",
            "",
            "Core Product Screens:",
        ] + [f"  ◈  {s}" for s in screens],
        0.55, 2.3, 5.4, 4.4, size=13, color=TEXT_BODY, font_name=font_name
    )

    # Right: Browser mockup graphic (Visual UI component!)
    mock_x, mock_y, mock_w, mock_h = 6.6, 1.6, 6.2, 5.4
    # Browser border
    _rect(slide, mock_x, mock_y, mock_w, mock_h, CARD_BG, border=True, border_color=_mix(DARK_BG, primary, 0.35))
    
    # Browser toolbar bar
    bar_h = 0.45
    _rect(slide, mock_x, mock_y, mock_w, bar_h, _mix(DARK_BG, primary, 0.25))
    
    # Mock browser control dots
    dot_y = mock_y + 0.15
    _rect(slide, mock_x + 0.18, dot_y, 0.12, 0.12, DANGER)
    _rect(slide, mock_x + 0.38, dot_y, 0.12, 0.12, ACCENT2)
    _rect(slide, mock_x + 0.58, dot_y, 0.12, 0.12, ACCENT)

    # Mock sidebar nav
    _rect(slide, mock_x, mock_y + bar_h, 1.6, mock_h - bar_h, _mix(DARK_BG, primary, 0.15))
    for i in range(len(screens[:4])):
        screen_y = mock_y + bar_h + 0.25 + (i * 0.45)
        # Highlight active item
        bg_col = _mix(DARK_BG, primary, 0.45) if i == 0 else _mix(DARK_BG, primary, 0.08)
        _rect(slide, mock_x + 0.1, screen_y, 1.4, 0.35, bg_col)
        # Use white text on active item
        lbl_col = WHITE if i == 0 else TEXT_BODY
        _txt(slide, screens[i][:15], mock_x + 0.2, screen_y + 0.05, 1.2, 0.25, size=9, color=lbl_col, font_name=font_name)

    # Mock main area content dashboard
    content_x = mock_x + 1.75
    _rect(slide, content_x, mock_y + bar_h + 0.2, mock_w - 1.9, 0.8, CARD_BG)
    _txt(slide, design_goal[:60] + "...", content_x + 0.15, mock_y + bar_h + 0.25, mock_w - 2.2, 0.7, size=10, italic=True, color=TEXT_MUTED, font_name=font_name)

    # Mini mock cards
    card_y = mock_y + bar_h + 1.2
    _rect(slide, content_x, card_y, 1.9, 1.4, CARD_BG)
    _txt(slide, "📊 Metrics Overview", content_x + 0.1, card_y + 0.1, 1.7, 0.3, size=9, bold=True, color=accent, font_name=font_name)
    _rect(slide, content_x + 0.1, card_y + 0.5, 1.7, 0.08, _mix(DARK_BG, primary, 0.25))
    _rect(slide, content_x + 0.1, card_y + 0.7, 1.4, 0.08, _mix(DARK_BG, primary, 0.25))
    _rect(slide, content_x + 0.1, card_y + 0.9, 1.1, 0.08, _mix(DARK_BG, primary, 0.25))

    _rect(slide, content_x + 2.1, card_y, 2.1, 1.4, CARD_BG)
    _txt(slide, "👤 Recent Actions", content_x + 2.2, card_y + 0.1, 1.9, 0.3, size=9, bold=True, color=primary, font_name=font_name)
    _rect(slide, content_x + 2.2, card_y + 0.5, 1.9, 0.08, _mix(DARK_BG, accent, 0.25))
    _rect(slide, content_x + 2.2, card_y + 0.7, 1.5, 0.08, _mix(DARK_BG, accent, 0.25))
    _rect(slide, content_x + 2.2, card_y + 0.9, 1.7, 0.08, _mix(DARK_BG, accent, 0.25))

    # Big table row list block
    table_y = mock_y + bar_h + 2.8
    _rect(slide, content_x, table_y, mock_w - 1.9, 1.7, CARD_BG)
    _txt(slide, "Primary Workspace Table", content_x + 0.15, table_y + 0.1, mock_w - 2.2, 0.3, size=10, bold=True, color=TEXT_TITLE, font_name=font_name)
    for j in range(3):
        row_y = table_y + 0.5 + (j * 0.35)
        _rect(slide, content_x + 0.15, row_y, mock_w - 2.2, 0.25, _mix(DARK_BG, primary, 0.12))
        _txt(slide, f"Row {j+1}", content_x + 0.25, row_y + 0.02, 1.0, 0.2, size=8, color=TEXT_BODY, font_name=font_name)
        _rect(slide, content_x + mock_w - 2.9, row_y + 0.05, 0.7, 0.15, accent) # Mock status badge


def _slide_tech_schema(
    prs: Presentation,
    title: str,
    stack_items: list[str],
    entities: list[dict[str, Any]],
    primary: RGBColor,
    accent: RGBColor,
    card_bg: RGBColor,
    font_name: str = "Outfit",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    _rect(slide, 0, 0, 0.1, 7.5, primary)
    _txt(slide, f"🏗  {title}", 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    # Left: Core stack list
    _rect(slide, 0.35, 1.6, 5.8, 5.4, card_bg)
    _txt(slide, "Technology Stack", 0.55, 1.8, 5.4, 0.4, size=16, bold=True, color=primary, font_name=font_name)
    _multiline_txt(slide, stack_items, 0.55, 2.3, 5.4, 4.4, size=14, color=TEXT_BODY, font_name=font_name)

    # Right: Database schema diagrams (visual component!)
    db_x = 6.6
    for i, ent in enumerate(entities[:2]):
        ent_y = 1.6 + (i * 2.7)
        ent_name = ent.get("name", f"Table_{i+1}")
        fields = ent.get("fields", ["id:uuid", "created_at:datetime"])[:4]

        # Table header
        _rect(slide, db_x, ent_y, 6.2, 0.5, _mix(DARK_BG, accent, 0.35))
        _txt(slide, f"📂  TABLE: {ent_name.upper()}", db_x + 0.25, ent_y + 0.1, 5.7, 0.4, size=14, bold=True, color=accent, font_name=font_name)

        # Table body / fields list
        _rect(slide, db_x, ent_y + 0.5, 6.2, 1.8, CARD_BG, border=True, border_color=_mix(DARK_BG, accent, 0.25))
        for j, field in enumerate(fields):
            field_y = ent_y + 0.6 + (j * 0.4)
            _txt(slide, f"🔑 {field}" if j==0 else f"   {field}", db_x + 0.25, field_y, 5.7, 0.35, size=11, color=TEXT_BODY, font_name=font_name)

    # Draw ERD relational vector line linking Table 1 and Table 2 (Visual architecture element!)
    if len(entities) > 1:
        # Vertical connector line
        _rect(slide, 9.7, 3.9, 0.04, 0.4, accent)
        # Dot nodes on ends
        _rect(slide, 9.67, 3.87, 0.1, 0.1, primary)
        _rect(slide, 9.67, 4.27, 0.1, 0.1, primary)


def _slide_score(
    prs: Presentation,
    score: dict[str, Any],
    primary: RGBColor,
    accent: RGBColor,
    card_bg: RGBColor,
    font_name: str = "Outfit",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    # Slide Title - Use glowing Neon Cyan
    _rect(slide, 0, 0, 0.1, 7.5, accent)
    _txt(slide, "📊  AI Startup Score", 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=NEON_CYAN, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, accent)
    # Subtitle - Use glowing Neon Green
    _txt(slide, "Multi-agent seed consensus — rated 1 to 10 across key dimensions",
         0.35, 1.42, 12.5, 0.5, size=13, italic=True, color=NEON_GREEN, font_name=font_name)

    if not score:
        _txt(slide, "Score data not available", 0.5, 3.0, 12.0, 1.0, size=20, color=TEXT_MUTED, align=PP_ALIGN.CENTER, font_name=font_name)
        return

    metrics = list(score.items())

    # Left: small layout cards (3 rows, 2 columns)
    cols = 2
    card_w = 2.85
    card_h = 1.6
    x_offs = [0.35, 3.4]
    y0 = 2.1
    row_gap = 1.85

    for i, (metric, val) in enumerate(metrics[:6]):
        col = i % cols
        row = i // cols
        x = x_offs[col]
        y = y0 + row * row_gap

        fval  = float(val)
        sc    = _score_color(fval, primary, accent) # Dynamic neon color

        # Card shape
        _rect(slide, x, y, card_w, card_h, card_bg)
        _rect(slide, x, y, card_w, 0.06, sc)

        # Score rating - Guaranteed to be bright neon
        _txt(slide, f"{fval:.0f}", x + 0.15, y + 0.1, 0.8, 0.8, size=38, bold=True, color=sc, font_name=font_name)
        _txt(slide, "/10", x + 0.75, y + 0.45, 0.6, 0.3, size=12, color=WHITE, font_name=font_name)

        # Metric label - Use pure high-contrast white inside dark cards
        label = metric.replace("_", " ").title()
        _txt(slide, label, x + 0.15, y + 1.1, card_w - 0.3, 0.45, size=11, bold=True, color=WHITE, font_name=font_name)

    # Right: Native PowerPoint Column Chart! (Visual graphic component)
    chart_data = CategoryChartData()
    chart_data.categories = [m.replace("_", " ").title() for m, _ in metrics[:6]]
    chart_data.add_series('Rating', tuple(float(val) for _, val in metrics[:6]))
    
    cx, cy, cw, ch = 6.6, 2.1, 6.2, 4.8
    # Add chart shape
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(cx), Inches(cy), Inches(cw), Inches(ch),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    
    # ── Style Chart Axis Labels and Title for dark mode readability ───────────
    chart.has_title = True
    title_tf = chart.chart_title.text_frame
    title_tf.text = "Metric Rating Analysis"
    if title_tf.paragraphs and title_tf.paragraphs[0].runs:
        title_run = title_tf.paragraphs[0].runs[0]
        title_run.font.color.rgb = NEON_CYAN
        title_run.font.size = Pt(13)
        title_run.font.bold = True
        title_run.font.name = font_name

    # X-Axis Category Labels (Feasibility, Innovation, etc.)
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.color.rgb = WHITE
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.name = font_name
    cat_axis.tick_labels.font.bold = True

    # Y-Axis Numerical Labels (0, 1, 2, 3...)
    val_axis = chart.value_axis
    val_axis.tick_labels.font.color.rgb = TEXT_BODY
    val_axis.tick_labels.font.size = Pt(9.5)
    val_axis.tick_labels.font.name = font_name

    # Color the columns with the brand primary color
    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = primary


def _slide_roadmap_timeline(
    prs: Presentation,
    title: str,
    roadmap: list[str],
    primary: RGBColor,
    accent: RGBColor,
    card_bg: RGBColor,
    font_name: str = "Outfit",
) -> None:
    slide = _blank_slide(prs)
    _fill_bg(slide, DARK_BG, SLIDE_BG_PATH)

    _rect(slide, 0, 0, 0.1, 7.5, primary)
    _txt(slide, f"🗺  {title}", 0.35, 0.3, 12.5, 0.9, size=30, bold=True, color=TEXT_TITLE, font_name=font_name)
    _rect(slide, 0.35, 1.3, 12.65, 0.055, primary)

    display_phases = roadmap[:3]
    if not display_phases:
        display_phases = ["Phase 1: Core MVP", "Phase 2: Growth", "Phase 3: Scale"]

    # Horizontal timeline axis arrow (representing future progression!)
    line_y = 3.6
    _draw_right_arrow(slide, 0.8, line_y - 0.06, 11.7, primary)

    # Calculate horizontal node placements
    total_nodes = len(display_phases)
    col_w = 3.6
    card_h = 1.8
    
    # Placements: Node 1 (Top), Node 2 (Bottom), Node 3 (Top) to create grid
    for i, phase in enumerate(display_phases):
        is_top = i % 2 == 0
        x = 0.8 + i * 4.0
        y = 1.6 if is_top else 4.2
        
        # Connect node to axis line via vertical connector pins
        conn_x = x + (col_w / 2.0)
        conn_y = 3.3 if is_top else 3.8
        conn_h = 0.35
        _rect(slide, conn_x - 0.03, conn_y, 0.06, conn_h, accent)

        # Timeline node card
        _rect(slide, x, y, col_w, card_h, card_bg)
        _rect(slide, x, y, col_w, 0.06, accent)

        # Phase label
        _txt(slide, f"PHASE {i+1}", x + 0.25, y + 0.15, col_w - 0.5, 0.35, size=11, bold=True, color=accent, font_name=font_name)

        # Phase text content
        _txt(slide, phase, x + 0.25, y + 0.55, col_w - 0.5, card_h - 0.7, size=12, color=TEXT_BODY, font_name=font_name)
